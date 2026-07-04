"""Generate the Phase 2 conference deck (CLV ranking -> uplift-aware targeting).
Same visual format as the Phase 1 deck (presentation/generate_conference_slides.py)."""
from __future__ import annotations
from pathlib import Path
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]          # phase2/
FIG = ROOT / "figures"
OUT = ROOT / "presentation" / "Group11_Phase2_Slides.pptx"
TMP = ROOT / "presentation" / "_charts"

NAVY = RGBColor(0x1A, 0x3A, 0x5C)
TEAL = RGBColor(0x00, 0x7A, 0x87)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
RED = RGBColor(0xC0, 0x39, 0x2B)


def _bar(path, labels, vals, ylabel, title, colors, ylim=None, fmt="{:.3f}", hline=None):
    TMP.mkdir(parents=True, exist_ok=True)
    fig, ax = plt.subplots(figsize=(7, 3.8))
    bars = ax.bar(labels, vals, color=colors, edgecolor="white", linewidth=0.8)
    ax.set_ylabel(ylabel)
    if ylim:
        ax.set_ylim(*ylim)
    if hline is not None:
        ax.axhline(hline, color="gray", ls="--", lw=0.8, alpha=0.6)
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width() / 2, v + (max(vals) - min(min(vals), 0)) * 0.02,
                fmt.format(v), ha="center", fontsize=9, fontweight="bold")
    ax.set_title(title, fontsize=11, fontweight="bold")
    plt.tight_layout()
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def chart_qini():
    p = TMP / "qini.png"
    return _bar(p, ["Response\nmodel", "Random", "T-Learner", "X-Learner", "S-Learner"],
                [-0.007, -0.002, 0.010, 0.012, 0.014], "Qini AUC",
                "Uplift ranking (RQ3): response < random",
                ["#e74c3c", "#95a5a6", "#5dade2", "#3498db", "#2E86AB"], ylim=(-0.02, 0.02))


def chart_profit():
    p = TMP / "profit.png"
    return _bar(p, ["Value\n(CLV)", "RFM", "Uplift", "Value-\nadjusted"],
                [-0.65, -0.49, 0.49, 0.74], "Profit @10% (M rubles, point est.)",
                "Targeting profit @10% (RQ4)",
                ["#e74c3c", "#f39c12", "#5dade2", "#2E86AB"], ylim=(-1.0, 1.1), fmt="{:+.2f}", hline=0)


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _bar_title(self, s, title, subtitle=""):
        bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.75))
        p = tb.text_frame.paragraphs[0]; p.text = title
        p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE
        if subtitle:
            sb = s.shapes.add_textbox(Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.4))
            sp = sb.text_frame.paragraphs[0]; sp.text = subtitle
            sp.font.size = Pt(14); sp.font.color.rgb = TEAL

    def _bullets(self, s, items, left=0.5, top=1.6, width=6.0, size=18):
        box = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
        tf = box.text_frame; tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = it; p.font.size = Pt(size); p.font.color.rgb = GRAY; p.space_after = Pt(9)

    def _img(self, s, path, left, top, width):
        if Path(path).exists():
            s.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))

    def _table(self, s, headers, rows, left=0.5, top=1.7, col_w=(2.8, 4.2)):
        nr, nc = len(rows) + 1, len(headers)
        t = s.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(sum(col_w)),
                               Inches(0.42 * nr)).table
        for j, h in enumerate(headers):
            c = t.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY
            for p in c.text_frame.paragraphs:
                p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = WHITE
        for i, row in enumerate(rows, 1):
            for j, v in enumerate(row):
                c = t.cell(i, j); c.text = v
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(11); p.font.color.rgb = GRAY

    def _kpis(self, s, kpis, top=1.7, w=2.8, gap=3.1):
        for i, (num, lab) in enumerate(kpis):
            left = 0.5 + i * gap
            box = s.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(1.5))
            box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.fill.background()
            tb = s.shapes.add_textbox(Inches(left), Inches(top + 0.15), Inches(w), Inches(1.2))
            tf = tb.text_frame
            p0 = tf.paragraphs[0]; p0.text = num; p0.font.size = Pt(24); p0.font.bold = True
            p0.font.color.rgb = WHITE; p0.alignment = PP_ALIGN.CENTER
            p1 = tf.add_paragraph(); p1.text = lab; p1.font.size = Pt(12)
            p1.font.color.rgb = RGBColor(0xBB, 0xDD, 0xEE); p1.alignment = PP_ALIGN.CENTER

    def _footer(self, s, text):
        box = s.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7))
        p = box.text_frame.paragraphs[0]; p.text = text
        p.font.size = Pt(13); p.font.italic = True; p.font.color.rgb = TEAL

    # ---------- slides ----------
    def s01(self):
        s = self._blank()
        bg = s.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(2.6))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ("From CLV Ranking to Uplift-Aware Targeting\n"
                  "Cross-Context Evidence from Online Retail, Grocery & Supermarket Campaign Data")
        p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
        meta = ["Vu Minh Anh Nguyen · Linh Hoang Nguyen (corresponding)",
                "FPT University · ADY201m — Phase 2",
                "CLV tells who is valuable · Uplift tells who is responsive · Value-adjusted uplift tells who is worth targeting"]
        mb = s.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.6))
        for i, line in enumerate(meta):
            mp = mb.text_frame.paragraphs[0] if i == 0 else mb.text_frame.add_paragraph()
            mp.text = line; mp.font.size = Pt(15); mp.font.color.rgb = RGBColor(0xBB, 0xDD, 0xEE)
            mp.alignment = PP_ALIGN.CENTER

    def s02(self):
        s = self._blank(); self._bar_title(s, "Agenda")
        self._bullets(s, [
            "1. Motivation — value vs persuadability (sure buyer vs persuadable)",
            "2. Research questions & three datasets with separated roles",
            "3. Methodology — Hurdle-Semantic CLV, uplift learners, value-adjusted policy",
            "4. Results — cross-context CLV, uplift > response, value-adjusted profit",
            "5. Honest significance, discussion & conclusion",
        ], top=1.6, size=20)
        self._footer(s, "Applied decision-support framework — integration + honest evaluation, not a new estimator.")

    def s03(self):
        s = self._blank(); self._bar_title(s, "Motivation: Value ≠ Persuadability")
        self._bullets(s, [
            "CLV ranks the most valuable customers — often already-loyal sure buyers.",
            "Uplift ranks the campaign-responsive — the persuadables.",
            "A high-value sure buyer returns anyway → wasted budget.",
            "A mid-value persuadable only buys when prompted → invisible to value models.",
            "Question: who is BOTH valuable AND incrementally persuadable?",
        ], width=6.0)
        self._img(s, FIG / "x5_policy_overlap.png", 7.0, 1.7, 5.8)
        self._footer(s, "Top-10% chosen by value vs by uplift overlap only 0.7% — near-disjoint customers.")

    def s04(self):
        s = self._blank(); self._bar_title(s, "Fit with ADY201m Requirements")
        rows = [
            ("Python pipeline", "Ingest → shared features → 14 models → evaluation"),
            ("SQL analysis", "DuckDB: RFM, segments, balance, overlap queries"),
            (">5 ML models", "8 CLV regressors + 6 uplift learners"),
            ("Dashboard / report", "Targeting dashboard + Springer PDF (12 pp.)"),
            ("AI Audit Log", "P2-01..P2-20, Human Delta + honest reporting"),
        ]
        self._table(s, ["Requirement", "Phase 2"], rows, top=1.7, col_w=(3.0, 6.5))
        self._footer(s, "Three datasets, deliberately separated roles — causal-clean design.")

    def s05(self):
        s = self._blank(); self._bar_title(s, "Research Questions")
        self._bullets(s, [
            "RQ1 — How robust is the Hurdle-Semantic CLV framework across e-commerce and grocery?",
            "RQ2 — Do behavioural, sequence, semantic features add value beyond RFM?",
            "RQ3 — Does uplift-aware targeting beat response- and value-based targeting?",
            "RQ4 — Does value-adjusted uplift (τ̂·V̂) improve budget-constrained efficiency?",
        ], top=1.7, size=19)
        self._footer(s, "RQ1–RQ2: CLV ranking | RQ3: uplift | RQ4: the fusion (central contribution).")

    def s06(self):
        s = self._blank(); self._bar_title(s, "Three Datasets, Separated Roles")
        self._kpis(s, [("0.8M", "Online Retail II tx"), ("2,500", "Dunnhumby hh"),
                       ("200,039", "X5 clients"), ("45.7M", "X5 purchase lines")])
        self._bullets(s, [
            "Online Retail II (UK e-commerce) — baseline CLV (read-only reuse).",
            "Dunnhumby (US grocery) — external CLV validation.",
            "X5 RetailHero (randomised SMS) — the only source of treatment-control uplift.",
            "Coupon redemption deliberately NOT used as treatment (post-campaign outcome).",
        ], top=3.5, width=12.0, size=17)

    def s07(self):
        s = self._blank(); self._bar_title(s, "Methodology — Four Modules")
        self._bullets(s, [
            "A. Hurdle CLV:  CLV̂ = p(return) × E[spend | return], lognormal correction.",
            "B. Semantic profiling:  PPMI co-occurrence → SVD product embeddings → taste vectors.",
            "C. Uplift learners:  S- / T- / X-Learner, class transformation vs response & random.",
            "D. Value-adjusted uplift (central):  Score = τ̂ᵢ · V̂ᵢ .",
        ], top=1.7, width=12.0, size=19)
        self._footer(s, "CLV = forecasting (walk-forward). Uplift = randomised cross-section (RCT).")

    def s08(self):
        s = self._blank(); self._bar_title(s, "Cross-Context Structure (RQ1)",
                                            "Two opposite retail regimes")
        self._img(s, FIG / "context_comparison.png", 0.6, 1.5, 12.1)
        self._footer(s, "E-commerce: 50% zero, skew 16.7 (heavy-tailed). Grocery: 7% zero, skew 2.5 (dense).")

    def s09(self):
        s = self._blank(); self._bar_title(s, "Cross-Context CLV Benchmark (RQ1, RQ2)",
                                            "Walk-forward, mean ± std — aligned with Phase 1")
        rows = [
            ("Online Retail II — Hurdle", "0.825 ± 0.05", "59.0%"),
            ("Online Retail II — Monetary", "0.822", "60.0%"),
            ("Dunnhumby — RandomForest", "0.862", "34.3%"),
            ("Dunnhumby — Hurdle", "0.855", "34.3%"),
        ]
        self._table(s, ["Model", "NG", "RC@10"], rows, top=1.8, col_w=(4.5, 2.2, 1.6))
        self._bullets(s, [
            "Hurdle competitive, not dominant — ties Monetary (p=0.55) / RF.",
            "Mechanism shifts: incidence (e-com) vs sequence/magnitude (grocery).",
            "Extra features > RFM significant in grocery (p=0.002), not e-com (p=0.17).",
            "Semantic: interpretation, not a ranking driver (ns both contexts).",
        ], left=8.6, top=1.8, width=4.4, size=14)

    def s10(self):
        s = self._blank(); self._bar_title(s, "Explainability — Stage-2 SHAP (RQ2)")
        self._img(s, FIG / "shap_crosscontext.png", 0.8, 1.7, 11.7)
        self._footer(s, "Monetary dominates spend in e-commerce; recent-sequence spend dominates in grocery.")

    def s11(self):
        s = self._blank(); self._bar_title(s, "Uplift-Aware Targeting (RQ3)",
                                            "Response model ranks BELOW random")
        self._img(s, chart_qini(), 0.6, 1.6, 6.5)
        self._bullets(s, [
            "Response model (\"who will buy\") Qini −0.007 — worse than random.",
            "Uplift learners positive: S-Learner Qini 0.014.",
            "Top-decile uplift ≈ 11% — ~3× the average effect.",
            "“Likely to buy” ≠ “persuadable” (paired bootstrap p < 0.001).",
        ], left=7.2, top=1.7, width=5.6, size=16)

    def s12(self):
        s = self._blank(); self._bar_title(s, "Value-Adjusted Targeting & Profit (RQ4)")
        self._img(s, chart_profit(), 0.6, 1.6, 6.5)
        self._bullets(s, [
            "Value & uplift top-10% overlap only 0.007 — different customers.",
            "Value-only can LOSE money at 10% (sure buyers, low uplift).",
            "Uplift-only maximises incremental conversions.",
            "Value-adjusted: highest point-estimate profit at tight budgets",
            "   (wins all 9 cost×margin settings; differences NOT significant).",
        ], left=7.2, top=1.7, width=5.6, size=15)

    def s13(self):
        s = self._blank(); self._bar_title(s, "Honest Significance — What Holds Up")
        rows = [
            ("Extra features > RFM (grocery)", "+0.044", "p < 0.01  ✓"),
            ("Extra features > RFM (e-com)", "+0.009", "p = 0.17  (ns)"),
            ("Semantic adds ranking", "~0.001", "ns (both)"),
            ("Sequence adds ranking (grocery)", "+0.021", "p < 0.01  ✓"),
            ("Uplift > response model", "—", "p < 0.001  ✓"),
            ("Policy dollar profit differences", "—", "not significant"),
        ]
        self._table(s, ["Claim", "ΔNG", "Evidence"], rows, top=1.7, col_w=(5.2, 1.6, 2.6))
        self._footer(s, "Reported transparently: ranking-level claims significant; dollar-level profit is directional.")

    def s14(self):
        s = self._blank(); self._bar_title(s, "Discussion — Transferable Insights")
        self._bullets(s, [
            "Do not target on CLV alone — selects sure buyers, can lose money.",
            "Do not target on response alone — can underperform random.",
            "Value-adjusted uplift suits budget-constrained settings.",
            "Explainability (SHAP) + calibrated uncertainty (CQR, 90%→90.2%/89.0%) build trust.",
            "0.7% overlap → a single objective leads to suboptimal / loss-making targeting.",
        ], width=12.0, size=18)
        self._footer(s, "Transferable principles for budget-constrained emerging markets — not numerical transfer.")

    def s15(self):
        s = self._blank(); self._bar_title(s, "Conclusion")
        self._bullets(s, [
            "Framework transfers but mechanism is context-dependent.",
            "Response-based targeting is worse than random for incremental campaigns.",
            "Value-adjusted uplift = most promising rule under tight budgets (honest caveats).",
            "Limitations: small grocery sample, single campaign, value proxy, modest signal.",
        ], width=12.0, size=18)
        box = s.shapes.add_shape(1, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.4))
        box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(11.9), Inches(1.05))
        p = tb.text_frame.paragraphs[0]
        p.text = ("CLV identifies who is valuable · uplift identifies who is responsive · "
                  "value-adjusted uplift identifies who is worth targeting")
        p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

    def build(self):
        for m in [self.s01, self.s02, self.s03, self.s04, self.s05, self.s06, self.s07,
                  self.s08, self.s09, self.s10, self.s11, self.s12, self.s13, self.s14, self.s15]:
            m()
        OUT.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(OUT))
        return OUT


if __name__ == "__main__":
    print("Saved:", Deck().build())
