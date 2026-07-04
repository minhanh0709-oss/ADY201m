# -*- coding: utf-8 -*-
"""Phase 2 slides in the OFFICIAL ADY201m course format.

Strategy: clone paper_final/ADY201m_Slide_Presentation_Final_Sample (2).pptx so the
theme / master / fonts / table styling / title bars are preserved exactly, then
replace text frames, table cells, and pictures in place with Phase 2 content.

Narrative (per user plan): CLV = who is valuable; Uplift = who is persuadable;
Value-adjusted uplift = who is worth targeting.

Output: phase2/presentation/Group11_Phase2_Slides_ADY201m.pptx  (NEW name; original untouched)
"""
import copy
import os
from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, "..", ".."))
SAMPLE = os.path.join(ROOT, "paper_final", "ADY201m_Slide_Presentation_Final_Sample (2).pptx")
FIG = os.path.join(HERE, "..", "figures")
SHOTS = os.path.join(ROOT, "dashboard", "screenshots")   # real dashboard screenshots
OUT = os.path.join(HERE, "Group11_Phase2_Slides_ADY201m_v4.pptx")


# ---------- helpers that PRESERVE formatting ----------
def _set_para_text(p, text):
    runs = p.runs
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p.add_run().text = text


def set_lines(tf, lines):
    """Fill a text frame with `lines`, reusing paragraph-level formatting
    (bullets, indent) by cloning the last existing paragraph as needed."""
    while len(tf.paragraphs) < len(lines):
        last = tf.paragraphs[-1]._p
        last.addnext(copy.deepcopy(last))
    while len(tf.paragraphs) > len(lines):
        tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
    for p, line in zip(tf.paragraphs, lines):
        _set_para_text(p, line)


def set_cell(cell, text):
    tf = cell.text_frame
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    _set_para_text(tf.paragraphs[0], text)


def set_table(tbl, rows):
    for r, row in enumerate(rows):
        if r >= len(tbl.rows):
            break
        for c, val in enumerate(row):
            if c < len(tbl.columns):
                set_cell(tbl.cell(r, c), val)


def replace_pics(slide, paths, base=FIG, fit=False):
    from PIL import Image as _Img
    pics = [sh for sh in slide.shapes if sh.shape_type == 13]
    for pic, path in zip(pics, paths):
        l, t, w, h = pic.left, pic.top, pic.width, pic.height
        pic._element.getparent().remove(pic._element)
        fp = os.path.join(base, path)
        if fit:  # preserve aspect ratio, centre inside the slot (no stretching)
            iw, ih = _Img.open(fp).size
            s = min(w / iw, h / ih)
            nw, nh = int(iw * s), int(ih * s)
            nl, nt = l + (w - nw) // 2, t + (h - nh) // 2
            slide.shapes.add_picture(fp, nl, nt, nw, nh)
        else:
            slide.shapes.add_picture(fp, l, t, w, h)


def tframes(slide):
    return [sh for sh in slide.shapes if sh.has_text_frame]


def tables(slide):
    return [sh for sh in slide.shapes if sh.has_table]


# ---------- build ----------
prs = Presentation(SAMPLE)
S = list(prs.slides)


# S1 — Title
t = tframes(S[0])
set_lines(t[0].text_frame, ["From CLV Ranking to Uplift-Aware Targeting: "
                            "A Cross-Context Study Across Three Retail Datasets"])
set_lines(t[1].text_frame, ["Ho Chi Minh City, 2026  —  Phase 2 moves from "
                            "“who is valuable?” to “who is worth targeting?”"])
set_lines(t[2].text_frame, ["Students:", "Nguyễn Vũ Minh Ánh - SE203412"])
set_lines(t[3].text_frame, ["Supervisor:", "Nguyễn Hoàng Linh"])

# S2 — Outline
set_lines(tframes(S[1])[1].text_frame, [
    "",
    "PART 1 — Problem & Data",
    "Problem Statement · Research Questions · Three datasets · SQL setup",
    "PART 2 — Analysis & Results",
    "CLV benchmark · Uplift learners · Value-adjusted policy",
    "PART 3 — Tool & AI Reflection",
    "Dashboard / report · AI Audit Log · Human Delta · Hallucination detection",
])

# S3 — Problem Statement (table 6x2) + footnote
set_table(tables(S[2])[0].table, [
    ["Evaluation Criteria", "Phase 2 — From Value to Worth-Targeting"],
    ["Describe the problem", "High-value customers are not always the best campaign targets; some "
     "high-CLV customers are already loyal and would purchase without incentives."],
    ["State the analysis objective", "Connect CLV ranking with uplift-aware targeting so the "
     "campaign reaches customers who are both valuable and persuadable."],
    ["Justify the approach", "CLV ranks who is valuable, uplift estimates who is persuadable, and "
     "value-adjusted uplift combines the two; all validated walk-forward."],
    ["Cite 1–2 reliable sources", "Ascarza (2018, JMR); Künzel et al. (2019, PNAS). Data: UCI "
     "Online Retail II, Dunnhumby Complete Journey, X5 RetailHero."],
    ["AI Audit Log required", "20 logged sessions (P2-01–P2-20) documenting prompts, validation, "
     "and corrections."],
])
set_lines(tframes(S[2])[1].text_frame, [
    "Core logic:  CLV ranks who is valuable  ·  Uplift estimates who is persuadable  ·  "
    "Value-adjusted uplift targets customers who are both."])

# S4 — Research Questions
set_lines(tframes(S[3])[1].text_frame, [
    "Research Questions (each linked to a SQL query / Python analysis):",
    "[RQ1] Does the Hurdle + semantic CLV framework remain competitive across e-commerce and "
    "grocery?",
    "           Link to: CLV walk-forward benchmark; SQL RFM and window tables.",
    "[RQ2] Which feature groups improve ranking — RFM, sequence, or semantic?",
    "           Link to: leave-one-group-out feature ablation.",
    "[RQ3] Can uplift models identify persuadable customers better than response models?",
    "           Link to: X5 uplift benchmark (Qini / AUUC / uplift@K).",
    "[RQ4] Does value-adjusted uplift improve held-out targeting profit?",
    "           Link to: held-out policy profit simulation.",
])

# S5 — Data Understanding (table 7x2)
set_table(tables(S[4])[0].table, [
    ["What to present", "Phase 2 — Three Datasets, Three Roles"],
    ["Data source", "Online Retail II (UK e-commerce); Dunnhumby Complete Journey (US grocery); "
     "X5 RetailHero (randomised SMS campaign)."],
    ["Dataset size", "Online Retail II ≈ 5,852 customers; Dunnhumby ≈ 2,500 households; "
     "X5 ≈ 200k customers (held-out 60,010)."],
    ["Dataset roles", "Online Retail II = baseline CLV; Dunnhumby = external CLV validation; "
     "X5 = treatment–control uplift."],
    ["Key variables", "Future 6-month CLV (ranking); treatment flag + conversion (uplift); "
     "RFM, recent sequence, and semantic features."],
    ["Important design", "Dunnhumby coupon redemption is NOT used as treatment; X5 is the only "
     "uplift dataset, keeping the causal scope clean."],
    ["AI Audit Log", "Dataset-role and split-design prompts logged (P2-01–P2-06)."],
])

# S6 — EDA / Data Gap (3 captions + 3 images)
t = tframes(S[5])
set_lines(t[1].text_frame, ["Cross-context: e-commerce is sparse & zero-inflated, grocery is denser"])
set_lines(t[2].text_frame, ["Grocery CLV is lighter-tailed and more sequence-driven"])
set_lines(t[3].text_frame, ["X5: treatment vs control shows a positive raw uplift (+3.3pp)"])
replace_pics(S[5], ["context_comparison.png", "audit_dunnhumby_zero_skew.png",
                    "audit_x5_uplift_signal.png"])

# S7 — SQL Overview
set_lines(tframes(S[6])[1].text_frame, [
    "SQL (DuckDB) builds transparent, reproducible analytical tables before Python modeling:",
    "•  RFM and historical value (V) tables per dataset.",
    "•  Walk-forward observation / prediction windows — past data predicts future value.",
    "•  X5 treatment / control cohorts for uplift.",
    "•  Product / category aggregation.",
    "•  Dashboard-ready output tables.",
    "Key message: SQL supports reproducible business analysis before Python modeling.",
])

# S8 — Representative SQL Queries (table 4x3)
set_table(tables(S[7])[0].table, [
    ["Query → RQ", "Representative SQL (DuckDB)", "Insight"],
    ["Q1 · Customer RFM / value proxy → RQ1/RQ4",
     "SELECT CustomerID, MAX(date), COUNT(*), SUM(spend) FROM tx GROUP BY CustomerID",
     "Customer-level historical behavior and the value proxy V."],
    ["Q2 · Walk-forward CLV windows → RQ1/RQ2",
     "WITH past AS (… < cutoff), fut AS (… >= cutoff) SELECT … JOIN ON CustomerID",
     "Use past data to predict future value, with no leakage."],
    ["Q3 · X5 treatment–control cohorts → RQ3/RQ4",
     "SELECT client, treatment_flag, MAX(converted) FROM x5 GROUP BY client",
     "Compare treated versus control customers for uplift."],
])

# S9 — Python Pipeline (5 step label/desc pairs + footnote)
t = tframes(S[8])
pairs = [
    ("1 · Clean & Validate",
     "Drop missing IDs, keep Quantity > 0 and Price > 0, remove cancellations and duplicates"),
    ("2 · Build Features",
     "RFM, behavioral, recent-sequence, and semantic taste features for every dataset"),
    ("3 · CLV Walk-forward Windows",
     "Past-only observation windows predict the future 6-month CLV, with no leakage"),
    ("4 · Train CLV & Uplift Models",
     "Two-stage Hurdle for value; S / T / X-Learner meta-learners for persuadability"),
    ("5 · Evaluate",
     "Ranking (NG, RC@10%), uplift (Qini, AUUC), profit@K, SHAP drivers, and CQR intervals"),
]
for k, (lab, desc) in enumerate(pairs):
    set_lines(t[1 + 2 * k].text_frame, [lab])
    set_lines(t[2 + 2 * k].text_frame, [desc])
set_lines(t[11].text_frame, ["Key message: Python is the full modeling and evaluation pipeline."])

# S10 — Analysis with Python (table 7x2)
set_table(tables(S[9])[0].table, [
    ["What to present", "Phase 2 detail"],
    ["CLV analysis", "Two-stage Hurdle model, RFM baseline, and tree models under walk-forward "
     "validation."],
    ["Uplift analysis", "S-Learner, T-Learner, and X-Learner versus a response model and a random "
     "baseline."],
    ["Metrics", "NG, RC@10%, Qini, AUUC, uplift@K, and profit@K."],
    ["Feature engineering", "RFM, monetary value proxy, recent sequence, and PPMI + SVD semantic "
     "embeddings."],
    ["Conclusion", "Different business questions require different metrics; ranking is not the "
     "same as incremental targeting."],
    ["AI Audit Log", "Modeling and evaluation prompts logged (P2-07–P2-16)."],
])
set_lines(tframes(S[9])[1].text_frame, [
    "Key message: different business questions require different metrics."])

# S11 — Visualization Insights (body + 3 images)
set_lines(tframes(S[10])[1].text_frame, [
    "•  Insight 1 — The response model performs worse than random for incremental targeting "
    "(Qini −0.0071).",
    "•  Insight 2 — Uplift learners recover persuadable customers (S-Learner Qini 0.0136, "
    "uplift@10% ≈ 11%).",
    "•  Insight 3 — Top-value and top-uplift customers almost do not overlap (0.007 at the top 10%).",
    "•  Insight 4 — Value-adjusted uplift gives the highest point-estimate profit under tight budgets.",
    "Key message: value and persuadability are different.",
])
replace_pics(S[10], ["x5_qini_curves.png", "x5_policy_overlap.png", "x5_policy_profit.png"])

# S12 — Regression / Uplift Overview
set_lines(tframes(S[11])[1].text_frame, [
    "CLV side — the Hurdle model separates purchase incidence from spending amount.",
    "•  Stage 1 estimates the return probability p; Stage 2 estimates positive spend m; CLV = p · m.",
    "Uplift side — meta-learners estimate the incremental campaign effect, "
    "uplift = E[Y | treated] − E[Y | control].",
    "Why both are needed: CLV tells expected value; uplift tells campaign responsiveness.",
])

# S13 — Model Comparison (table 7x2)
set_table(tables(S[12])[0].table, [
    ["What to present", "Phase 2 detail"],
    ["CLV models", "Monetary, RFM, Ridge, Random Forest, XGBoost, LightGBM, and the two-stage "
     "Hurdle."],
    ["Uplift models", "Random, Response model, S-Learner, T-Learner, and X-Learner."],
    ["Dependent & independent", "CLV: target = future 6-month value; features = RFM / sequence / "
     "semantic. Uplift: target = conversion under treatment."],
    ["Result charts", "Walk-forward NG and RC@10% per context; Qini curves; profit versus "
     "targeting depth."],
    ["Interpretation & limits", "The framework is compared against simple and strong baselines; "
     "Hurdle is competitive, not dominant (only 3–4 windows, so tests are diagnostics)."],
    ["AI Audit Log", "Benchmark and significance-testing prompts logged (P2-08–P2-13)."],
])

# S14 — Results Interpretation (body + 1 image)
set_lines(tframes(S[13])[1].text_frame, [
    "Main findings:",
    "•  Hurdle remains competitive, not dominant (NG ≈ 0.825 e-commerce, 0.855 grocery; gaps not "
    "significant).",
    "•  E-commerce ranking is incidence-driven; grocery is sequence-driven (drop-sequence "
    "ΔNG +0.021, p = 0.001).",
    "•  Semantic features mainly support interpretation, not significant lift (p ≥ 0.58).",
    "•  Value-adjusted uplift is useful for budget-constrained targeting; SHAP drivers shift by "
    "context (figure).",
])
replace_pics(S[13], ["shap_crosscontext.png"])

# S15 — Dashboard / Targeting Report (table 7x2)
set_table(tables(S[14])[0].table, [
    ["Requirement", "Phase 2 detail"],
    ["Accepted tool", "Interactive Plotly Dash app (5 pages) — Phase 2 view: the “Uplift "
     "Targeting (P2)” tab (route /uplift); the output is decision support, not only prediction."],
    ["Report outputs", "CLV ranking, uplift ranking, value-adjusted targeting list, and policy "
     "comparison."],
    ["Interpretation", "SHAP drivers per context and CQR uncertainty intervals."],
    ["KPI / Summary metrics", "Held-out profit, incremental conversions, Qini AUC, and "
     "Revenue Capture@10%."],
    ["Filter / Slicer function", "Targeting depth K (10 / 20 / 30%), margin m, and contact cost c."],
    ["AI Audit Log", "Report and figure-generation prompts logged (P2-14–P2-18)."],
])

# S16 — Dashboard Demo (4 captions + 4 REAL dashboard screenshots from the /uplift page)
t = tframes(S[15])
set_lines(t[1].text_frame, ["1 · Customer ranking (CLV + uplift)  —  live dashboard, “Uplift Targeting (P2)” tab (/uplift)"])
set_lines(t[2].text_frame, ["2 · Policy comparison  —  incremental revenue by policy (/uplift)"])
set_lines(t[3].text_frame, ["3 · Profit simulation  —  held-out profit by policy (/uplift)"])
set_lines(t[4].text_frame, ["4 · Explainability & uncertainty  —  SHAP drivers + CQR coverage (/uplift)"])
replace_pics(S[15], ["real_uplift_1_ranking.png", "real_uplift_2_policy.png",
                     "real_uplift_3_profit.png", "real_uplift_4_explain.png"],
             base=SHOTS, fit=True)

# S17 — AI Audit Log (4 representative entries)
set_lines(tframes(S[16])[1].text_frame, [
    "AI Audit Log — logged sessions with prompt · AI response · evaluation · Human Delta. "
    "Four representative entries:",
    "•  [E3] Dataset roles / causal-clean — Dunnhumby coupon redemption is a post-campaign "
    "outcome, so it is excluded as a treatment; uplift is restricted to the randomised X5.",
    "•  [E13] Uplift vs response — the response model ranks below random; “who is likely to "
    "buy” ≠ “who is persuadable”.",
    "•  [E17] Honest reporting — value-adjusted uplift has the highest point-estimate profit, "
    "but bootstrap CIs overlap, so it is reported as directional, not significant.",
    "•  [E18] Claim + format fix — AI-generated san() stripped backslashes (± rendered as "
    "“pm”); rewritten, and every claim hedged to the statistical evidence.",
    "Key message: AI accelerated the work; every output was verified, corrected, and logged.",
])

# S18 — Human Delta
set_lines(tframes(S[17])[1].text_frame, [
    "Decisions the team made independently of AI:",
    "•  Do not use Dunnhumby coupon redemption as treatment (it would introduce selection bias).",
    "•  Use X5 only for uplift — the single clean randomised campaign.",
    "•  Replace the random split with walk-forward validation (temporal validity).",
    "•  Report non-significant profit differences honestly (point estimate, not “wins”).",
    "•  Reframe semantic features as interpretability, not a guaranteed performance gain.",
])

# S19 — Hallucination Detection
set_lines(tframes(S[18])[1].text_frame, [
    "AI mistakes that were caught and corrected with evidence:",
    "•  Overclaiming the CLV framework “transfers” → “remains competitive; mechanism shifts.”",
    "•  Overclaiming value-adjusted uplift “wins” → “highest point estimate, not significant.”",
    "•  Confusing response prediction with uplift → the response model ranks below random.",
    "•  Reusing Phase 1 numbers in Phase 2 → recomputed walk-forward, replaced an inflated NG ≈ 0.92.",
    "•  LaTeX / table formatting errors (the sanitiser rendered ± as “pm”) → the regex was rewritten.",
    "Key message: the project shows evidence-based correction, not blind AI use.",
])

# S20 — Closing / Q&A
set_lines(tframes(S[19])[1].text_frame, [
    "Three takeaways:",
    "•  CLV identifies who is valuable.",
    "•  Uplift identifies who is responsive.",
    "•  Value-adjusted uplift identifies who is worth targeting.",
    "Final message: marketing should target customers who are both valuable and incrementally "
    "persuadable — not only high-value customers.",
    "Q&A — thank you.",
])

prs.save(OUT)
print("saved:", OUT, "| slides:", len(S))
